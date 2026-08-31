"""
Generate slides.pptx — VoxTell & nnInteractive optimization deck.
Run: python make_slides.py

Design: one accent colour on a near-white ground. Whitespace separates sections
rather than full-width rules, which read as generated. No em dashes.
Every figure traces to a measured source — see SPEAKER_NOTES.md.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette: ground, ink, one accent ──────────────────────────────────────────
BG     = RGBColor(0xFD, 0xFD, 0xFD)   # near-white, neutral
INK    = RGBColor(0x1C, 0x20, 0x26)
ACCENT = RGBColor(0xA8, 0x6D, 0x18)
MUTED  = RGBColor(0x8C, 0x8F, 0x96)
HAIR   = RGBColor(0xE8, 0xE8, 0xE6)

L, CW = 0.95, 11.45

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = BG
    return s


def tx(s, text, l, t, w, h, size=13, bold=False, color=INK,
       align=PP_ALIGN.LEFT, space=None):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space:
            p.space_after = Pt(space)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = "Calibri"
        r.font.color.rgb = color
    return box


def rule(s, l, t, w, color=HAIR, thick=1):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(thick / 72))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def head(s, title, num):
    """Title, page number, and a short accent mark. No full-width bar."""
    tx(s, title, L, 0.6, 9.8, 0.6, size=27, bold=True)
    tx(s, num, 11.0, 0.77, 1.4, 0.4, size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    rule(s, L, 1.32, 0.9, ACCENT, thick=2.5)


def stat(s, value, label, l, t, w, vsize=46, color=INK):
    tx(s, value, l, t, w, 0.8, size=vsize, bold=True, color=color)
    tx(s, label, l, t + 0.82, w, 0.5, size=10, color=MUTED)


def bullets(s, items, l, t, w, size=13, gap=0.44, color=INK):
    for i, it in enumerate(items):
        y = t + i * gap
        sq = s.shapes.add_shape(1, Inches(l), Inches(y + 0.085),
                                Inches(0.07), Inches(0.07))
        sq.fill.solid()
        sq.fill.fore_color.rgb = ACCENT
        sq.line.fill.background()
        tx(s, it, l + 0.24, y, w - 0.24, gap, size=size, color=color)


def table(s, headers, rows, l, t, w, h, widths=None):
    tbl = s.shapes.add_table(len(rows) + 1, len(headers), Inches(l), Inches(t),
                             Inches(w), Inches(h)).table
    if widths:
        for i, ww in enumerate(widths):
            tbl.columns[i].width = Inches(ww)
    for c, htxt in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG
        r = cell.text_frame.paragraphs[0].add_run()
        r.text = htxt
        r.font.size = Pt(10); r.font.bold = True
        r.font.name = "Calibri"; r.font.color.rgb = MUTED
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            r = cell.text_frame.paragraphs[0].add_run()
            r.text = str(val)
            r.font.size = Pt(11)
            r.font.name = "Calibri"; r.font.color.rgb = INK
            r.font.bold = (ci == 0)
    return tbl


def source(s, text):
    tx(s, text, L, 6.9, CW, 0.4, size=9, color=MUTED)


# ═══ 1 — Title ═════════════════════════════════════════════════════════════
s = slide()
rule(s, L, 1.55, 1.1, ACCENT, thick=3)
tx(s, "CVPR 2025  ·  Medical Image Segmentation", L, 1.85, 9, 0.4,
   size=11, color=ACCENT)
tx(s, "Making Two Segmentation\nModels Faster", L, 2.35, 11, 1.9,
   size=44, bold=True, color=INK)
tx(s, "VoxTell  ·  nnInteractive  ·  NVIDIA H100 MIG", L, 4.35, 10, 0.4,
   size=15, color=MUTED)
tx(s, "Brian Xiao   ·   Fir cluster, Alliance Canada   ·   August 2026",
   L, 6.6, 10, 0.4, size=10, color=MUTED)

# ═══ 2 — The two models ════════════════════════════════════════════════════
s = slide()
head(s, "Two Models, Two Prompting Styles", "02")

tx(s, "VoxTell", L, 1.85, 5, 0.45, size=20, bold=True, color=ACCENT)
bullets(s, [
    "Our lab's CVPR submission",
    "Prompt is free text: “the spleen”",
    "Needs a 4B-parameter LLM to encode it first",
], L, 2.4, 5.2)

tx(s, "nnInteractive", 6.95, 1.85, 5, 0.45, size=20, bold=True, color=INK)
bullets(s, [
    "The CVPR challenge baseline",
    "Prompt is a 3-D bounding box",
    "No text encoder, far cheaper per prompt",
], 6.95, 2.4, 5.2)

tx(s, "Both already accurate.", L, 4.55, 10, 0.5, size=21, bold=True)
tx(s, "Can they be made faster without giving that up?",
   L, 5.15, 10, 0.5, size=17, color=MUTED)
source(s, "DSC is the Dice Similarity Coefficient: overlap between prediction and expert annotation, 0 to 1. Above roughly 0.7 is clinically usable.")

# ═══ 3 — Method + audit ════════════════════════════════════════════════════
s = slide()
head(s, "How the Number Was Measured", "03")

tx(s, "VoxTell speedup, as measurement error was removed", L, 1.7, 8, 0.35,
   size=11, color=MUTED)

for i, (val, why, col) in enumerate([
    ("26×",   "baseline ran on CPU",         MUTED),
    ("17.6×", "no warm-up, cache mismatch",  MUTED),
    ("7.1×",  "first arm ate start-up cost", MUTED),
    ("2.7×",  "measured correctly",          ACCENT),
]):
    x = L + i * 2.9
    tx(s, val, x, 2.15, 2.7, 0.85, size=42, bold=True, color=col)
    tx(s, why, x, 3.08, 2.7, 0.6, size=10, color=MUTED)

tx(s, "What Changed Each Time", L, 4.05, 6, 0.35, size=12, bold=True, color=ACCENT)
bullets(s, [
    "Hold precision constant, so quantization cannot pose as algorithmic gain",
    "Warm the GPU, text backbone and sliding-window path before timing",
    "Assert the embedding cache is empty before every cold measurement",
    "Repeat to n≥4 and quote the range, never a single run",
], L, 4.5, 11.3, size=13, gap=0.42)

tx(s, "No correction changed the code. Only how it was measured.",
   L, 6.35, 11, 0.4, size=15, bold=True)
source(s, "The largest error surfaced by running the same script on two GPUs: identical phases behaved differently on an RTX 4070 SUPER and an H100.")

# ═══ 4 — What made it faster ═══════════════════════════════════════════════
s = slide()
head(s, "VoxTell: What Made It Faster", "04")

table(s, ["Change", "Effect", "DSC"], [
    ("Sliding window", "tile_step 0.75 plus crop to non-zero, 25 to 9 patches", "+0.0003"),
    ("Embedding cache", "repeat prompts return a stored tensor", "identical"),
    ("INT4 backbone", "1.5× faster encode, 2 GB VRAM instead of 8 GB", "0.97 agreement"),
    ("Numba preprocess", "no measurable gain at this volume size", "unchanged"),
], L, 1.75, CW, 2.3, widths=[2.7, 6.55, 2.2])

stat(s, "2.7× mean", "abdominal CT, H100 MIG  ·  2.6 to 2.8× across 4 runs",
     L, 4.55, 6, color=ACCENT)
tx(s, "One run: 3.27s → 1.28s", 7.6, 4.7, 5, 0.6, size=20, color=MUTED)
source(s, "Case CT_AMOS_amos_0018 (63×512×512). Both arms run INT4, so precision is held constant and this is algorithmic gain only.")

# ═══ 5 — VoxTell accuracy ══════════════════════════════════════════════════
s = slide()
head(s, "VoxTell: Accuracy", "05")

stat(s, "+0.0003", "mean DSC change  ·  0.8090 → 0.8093", L, 1.9, 5.5, vsize=50)
tx(s, "65 objects across 5 abdominal CT cases", L, 3.35, 5.5, 0.4,
   size=13, color=MUTED)
bullets(s, [
    "Speed optimizations cost no accuracy",
    "Measured against expert annotation",
    "Same checkpoint and fold on both arms",
], L, 4.0, 5.4, size=12, gap=0.42)

tx(s, "One caveat, stated plainly", 7.0, 1.9, 4.8, 0.4, size=15, bold=True, color=ACCENT)
bullets(s, [
    "INT4 quantization is on by default",
    "0.97 agreement with full precision",
    "Segments 5.5% fewer voxels",
    "One-sided, so bias rather than noise",
    "Measured on one case, not the full set",
], 7.0, 2.45, 5.2, size=12, gap=0.42)
source(s, "VoxTell DSC from accuracy_results.csv. The INT4 comparison is n=1 and measures output agreement, not accuracy against ground truth.")

# ═══ 6 — nnInteractive result ══════════════════════════════════════════════
s = slide()
head(s, "nnInteractive: Compiling the Network", "06")

tx(s, "torch.compile(network, mode='reduce-overhead')", L, 1.75, 9, 0.4,
   size=15, color=ACCENT)
bullets(s, [
    "One line changed",
    "Fuses kernels and cuts per-call dispatch overhead",
], L, 2.3, 9, size=13)

stat(s, "1.33×", "per object, mean of 4 runs", L, 3.4, 5, color=ACCENT)
bullets(s, [
    "0.288s → 0.215s per object",
    "Range 1.28 to 1.39×, so a 33% gain against an 8% spread",
], L, 4.85, 5.6, size=12, gap=0.42)

stat(s, "+0.0002", "mean DSC change  ·  294 objects", 7.0, 3.4, 5)
bullets(s, [
    "No run showed degradation",
    "Speed and accuracy from the same jobs",
], 7.0, 4.85, 5.2, size=12, gap=0.42)

tx(s, "CPU time held at 6:18 to 6:33 across runs while walltime fell, "
      "which is what a GPU-bound workload looks like.",
   L, 5.95, 11.4, 0.5, size=13)
source(s, "20 CT cases from the CVPR validation set, fold='all' checkpoint, H100 MIG 3g.40gb. CPU efficiency 13 to 20% of 8 cores; allocation since reduced to 2.")

# ═══ 7 — Break-even ════════════════════════════════════════════════════════
s = slide()
head(s, "The Speedup Is Not Free at the Start", "07")

tx(s, "Compiling costs time before the first prediction.", L, 1.8, 10, 0.5, size=18)

for i, (v, lab) in enumerate([
    ("23.6s", "one-time compile"),
    ("0.071s", "saved per object"),
    ("~22 cases", "to break even"),
]):
    stat(s, v, lab, L + i * 3.7, 2.6, 3.4,
         color=ACCENT if i == 2 else INK)

bullets(s, [
    "Batch of 881 validation cases: clearly worth it",
    "Radiologist with 3 scans: never recovered",
    "A batch optimization, and it should not be sold as anything else",
], L, 4.5, 11.3, size=15, gap=0.48)
source(s, "23.6s measured on node-local /tmp. The shared filesystem is slower, so treat it as a lower bound. About 331 objects at 14.7 objects per case.")

# ═══ 8 — Summary ═══════════════════════════════════════════════════════════
s = slide()
head(s, "End Results", "08")

table(s, ["", "Speedup", "Accuracy", "Evidence"], [
    ("VoxTell", "2.7×  (2.6 to 2.8×)", "+0.0003 DSC", "4 runs, abdominal CT"),
    ("nnInteractive", "1.33×  (1.28 to 1.39×)", "+0.0002 DSC", "4 runs, 294 objects"),
], L, 1.85, CW, 1.5, widths=[2.6, 3.0, 2.6, 3.25])

tx(s, "Still Open", L, 3.9, 5, 0.4, size=14, bold=True, color=ACCENT)
bullets(s, [
    "INT4 under-segments by 5.5% on one case",
    "Not yet measured across the validation set",
], L, 4.4, 5.9, size=12, gap=0.42, color=MUTED)

tx(s, "Next", 7.0, 3.9, 5, 0.4, size=14, bold=True, color=ACCENT)
bullets(s, [
    "Run INT4 against all 881 validation cases",
    "Turns a direction into a bound worth acting on",
], 7.0, 4.4, 5.2, size=12, gap=0.42, color=MUTED)

tx(s, "The most useful thing I built was a benchmark that kept catching itself.",
   L, 5.7, 11.5, 0.5, size=17, bold=True)

prs.save("slides.pptx")
print("Saved: slides.pptx  (8 slides)")
